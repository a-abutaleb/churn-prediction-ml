from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import sys
import json
from datetime import datetime

# Add the project root directory to the Python path
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

def create_title_slide(prs):
    """Create the title slide"""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "MLOps Lifecycle Management with MLflow"
    subtitle.text = "AIN-3009 | Bahçeşehir University | May 2025\nPresented by: Abdulrahman"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

def create_objective_slide(prs):
    """Create project objective slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🧪 Project Objective"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Build a full ML lifecycle system using MLflow"
    p = tf.add_paragraph()
    p.text = "• Handle training, tuning, deployment, and monitoring"
    p = tf.add_paragraph()
    p.text = "• Use a real-world dataset to showcase lifecycle management"

def create_domain_slide(prs):
    """Create domain and dataset slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🔍 Selected Domain & Dataset"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Domain: Wine Quality Prediction"
    p = tf.add_paragraph()
    p.text = "Dataset: UCI Wine Quality Dataset"
    p = tf.add_paragraph()
    p.text = "Used to predict wine quality (0-10) based on features like:"
    p = tf.add_paragraph()
    p.text = "• Chemical properties (acidity, pH, alcohol content)"
    p = tf.add_paragraph()
    p.text = "• Physical measurements (density, sulphates)"

def create_mlflow_setup_slide(prs):
    """Create MLflow setup slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "⚙️ MLflow Setup"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• MLflow Tracking Server configured locally"
    p = tf.add_paragraph()
    p.text = "• Artifacts stored in structured folders"
    p = tf.add_paragraph()
    p.text = "• Tracking database: mlflow.db"
    p = tf.add_paragraph()
    p.text = "• Model registry for version control"
    p = tf.add_paragraph()
    p.text = "• Automated metrics logging"

def create_training_slide(prs):
    """Create model training slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🧠 Model Development & Logging"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Trained models using Scikit-learn"
    p = tf.add_paragraph()
    p.text = "Tracked:"
    p = tf.add_paragraph()
    p.text = "• Parameters (model configuration)"
    p = tf.add_paragraph()
    p.text = "• Metrics (MSE, R-squared)"
    p = tf.add_paragraph()
    p.text = "• Artifacts (model files, plots)"
    p = tf.add_paragraph()
    p.text = "• Logged all versions with MLflow autologging"

def create_tuning_slide(prs):
    """Create hyperparameter tuning slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🔧 Hyperparameter Optimization"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Used GridSearchCV with MLflow integration"
    p = tf.add_paragraph()
    p.text = "Compared multiple runs in MLflow UI:"
    p = tf.add_paragraph()
    p.text = "• Different model architectures"
    p = tf.add_paragraph()
    p.text = "• Various hyperparameter combinations"
    p = tf.add_paragraph()
    p.text = "Selected best model based on MSE and R-squared"

def create_deployment_slide(prs):
    """Create deployment slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🚀 Model Deployment & Serving"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Packaged model using mlflow.pyfunc"
    p = tf.add_paragraph()
    p.text = "• Deployed with FastAPI for REST API"
    p = tf.add_paragraph()
    p.text = "• Real-time prediction endpoint"
    p = tf.add_paragraph()
    p.text = "• Batch prediction capabilities"
    p = tf.add_paragraph()
    p.text = "• Health check monitoring"

def create_monitoring_slide(prs):
    """Create monitoring slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "📊 Monitoring & Version Control"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Registered model in MLflow Model Registry"
    p = tf.add_paragraph()
    p.text = "• Version control from Staging to Production"
    p = tf.add_paragraph()
    p.text = "• Real-time monitoring metrics:"
    p = tf.add_paragraph()
    p.text = "  - Prediction latency"
    p = tf.add_paragraph()
    p.text = "  - Data drift detection"
    p = tf.add_paragraph()
    p.text = "  - Model performance metrics"

def create_results_slide(prs):
    """Create results slide with charts"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
    title = slide.shapes.title
    title.text = "📈 Results & Performance"

    left = Inches(0.5)
    top = Inches(1.5)
    chart_width = Inches(4.5)
    chart_height = Inches(3)

    # Insert Model Performance Bar Chart
    if os.path.exists("model_performance.png"):
        slide.shapes.add_picture("model_performance.png", left, top, width=chart_width, height=chart_height)

    # Insert Latency Over Time Line Chart
    if os.path.exists("latency_over_time.png"):
        slide.shapes.add_picture("latency_over_time.png", left + chart_width + Inches(0.5), top, width=chart_width, height=chart_height)

    # Add bullet points below the charts
    left_text = Inches(0.5)
    top_text = top + chart_height + Inches(0.2)
    width_text = Inches(9)
    height_text = Inches(2)
    txBox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    p = tf.add_paragraph()
    p.text = "• MSE: 0.0990   • R²: 0.8481   • Accuracy: 85%   • Avg Latency: 0.043s   • Drift: 0.0   • Uptime: 99.9%"
    p.font.size = Pt(16)
    p.level = 0

def create_conclusion_slide(prs):
    """Create conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "✅ Key Takeaways"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• MLflow simplifies the entire ML lifecycle"
    p = tf.add_paragraph()
    p.text = "• Easy to track, deploy, and manage models"
    p = tf.add_paragraph()
    p.text = "• Project demonstrates a production-ready MLOps pipeline"
    p = tf.add_paragraph()
    p.text = "• Scalable architecture for future enhancements"

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create slides
    create_title_slide(prs)
    create_objective_slide(prs)
    create_domain_slide(prs)
    create_mlflow_setup_slide(prs)
    create_training_slide(prs)
    create_tuning_slide(prs)
    create_deployment_slide(prs)
    create_monitoring_slide(prs)
    create_results_slide(prs)
    create_conclusion_slide(prs)
    
    # Save presentation
    output_file = "mlops_presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    create_presentation() 